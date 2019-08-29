# -*- coding: utf-8 -*-

"""

@author: kevic

"""

##Import modules and set up
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import random
import os.path


class color:   
   BOLD = '\033[1m'
   END = '\033[0m'
   
   

##Update this section ONLY each time
prs = Presentation(r'file location and presentation name')  #must fill out at start
image_folder = r'file locaton for images to be saved'  #must fill out at start
pres_name  = 'presentation name' 
get_photos = 'y'    #write 'y' if you want to save all the photos in a folder otherwise write 'n'
scan_text = 'y'     #write 'y' if you want to scan the text for a phrase otherwise write 'n'
searchText = 'Write Text Here' 
 



##Read the text
if scan_text.lower() == 'y':  #forcing lower case just in case
    textContent = []
     
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    textContent.append(run.text)   
    
###Search the text for strings or phrases
    textContent_lwr = []     #find all matches and place them in a list
    for x in textContent:
        textContent_lwr.append(x.lower()) 
    
    match = []               #make lower case since case sensitive
    for x in textContent_lwr:
        if searchText.lower() in x: 
            match.append(x)
    
    match_unq = []           #shorten list by only showing unique matches
    for x in match:
        if x not in match_unq:
            match_unq.append(x)
            
    print('Looking for paragraphs with the following string: ' + searchText + '.\n')
    
    if not match_unq:       #check if list is empty
        print('There was no match.')
    else: 
        print(color.BOLD + 'Match identified! See below.' + color.END)
        print(*match_unq,sep='\n')
    
    
    
##Save all pictures with presentation name in a folder
if get_photos.lower() == 'y':       #forcing lower case just in case

    def iter_picture_shapes(prs):
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    yield shape
    
    for picture in iter_picture_shapes(prs):
        
        image = picture.image    
        image_bytes = image.blob    #make binary data so can save wb
        
        image_name = str(random.randint(1,100000)) + pres_name + '_img.jpg'
        image_filename = os.path.join(image_folder, image_name)
        with open(image_filename, 'wb') as f:
           f.write(image_bytes)     #with write binary python will overwrite if the file exists 
                                    #although unlikely since a random number is always generated 
    
    print('\nPictures were saved in ' + image_folder + ' when this script ran.')
   
else:
    print('No images were found or no scan was conducted.')
    
if scan_text.lower() == 'n' and get_photos.lower() == 'n':  #forcing lower case just in case
    print('Nothing happened when you ran the code. Please change the parameters.')
    

        

        
        